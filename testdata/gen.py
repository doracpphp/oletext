#!/usr/bin/env python3
"""Generate every test fixture for oletext, in one script.

There are two generation backends, selected by subcommand:

OOXML samples (.docx/.xlsx/.pptx) use the pure-Python libraries:

    py -m pip install python-docx openpyxl python-pptx
    py testdata/gen.py samples     # sample.docx/.xlsx/.pptx  (feature coverage)
    py testdata/gen.py large       # big.docx/.xlsx/.pptx     (perf/bench)

Legacy binary samples (.doc/.xls/.ppt) use LibreOffice via UNO. Start a
headless listener first, then run with LibreOffice's bundled python:

    soffice --headless --invisible --nologo \
            "--accept=socket,host=localhost,port=2003;urp;" &
    "<LibreOffice>/program/python" testdata/gen.py shapes   # shapes.doc/.xls/.ppt
    "<LibreOffice>/program/python" testdata/gen.py big       # big.doc/.xls/.ppt

The asserted strings are checked by the *_test.go files, so keep them in
sync. testdata/ is git-ignored; only this script is committed.
"""

import os
import sys

OUT = os.path.dirname(os.path.abspath(__file__))


# ---------------------------------------------------------------------------
# OOXML samples (python-docx / openpyxl / python-pptx)
# ---------------------------------------------------------------------------

def make_docx():
    import docx
    from docx.oxml import parse_xml
    d = docx.Document()
    cp = d.core_properties
    cp.title, cp.author, cp.keywords = "Doc Property Title", "Doc Author", "alpha beta"
    d.add_heading("Report Title", level=1)
    d.add_paragraph("First body paragraph with plain text.")
    d.add_paragraph("日本語の段落テキスト。")
    table = d.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "CellA"
    table.rows[0].cells[1].text = "CellB"
    p = d.add_paragraph("This sentence is annotated.")
    d.add_comment(p.runs[0], text="A reviewer comment.", author="Reviewer")
    # A VML text box (shape) appended as raw WordprocessingML.
    textbox = parse_xml(
        '<w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
        ' xmlns:v="urn:schemas-microsoft-com:vml">'
        "<w:r><w:pict>"
        '<v:shape type="#_x0000_t202" style="width:200pt;height:40pt">'
        "<v:textbox><w:txbxContent>"
        "<w:p><w:r><w:t>Text inside a shape.</w:t></w:r></w:p>"
        "</w:txbxContent></v:textbox></v:shape>"
        "</w:pict></w:r></w:p>"
    )
    d.element.body.append(textbox)
    d.save(os.path.join(OUT, "sample.docx"))


def make_xlsx():
    import openpyxl
    from openpyxl.comments import Comment
    wb = openpyxl.Workbook()
    wb.properties.title, wb.properties.creator = "Sheet Property Title", "Sheet Author"
    ws1 = wb.active
    ws1.title = "Numbers"
    ws1["A1"], ws1["B1"] = "Label", "Value"
    ws1["A2"], ws1["B2"] = "Answer", 42
    ws1["A3"] = "日本語セル"
    ws1["B2"].comment = Comment("A cell comment.", "Author")
    ws1.oddHeader.center.text = "Printed Report Header"
    ws1.oddFooter.right.text = "Sheet Footer Note"
    ws2 = wb.create_sheet("Second")
    ws2["A1"] = "On the second sheet"
    wb.save(os.path.join(OUT, "sample.xlsx"))


def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.core_properties.title, prs.core_properties.author = "Deck Property Title", "Deck Author"
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text_frame.text = "Slide one shape text."
    s1.notes_slide.notes_text_frame.text = "Speaker notes for slide one."
    s2 = prs.slides.add_slide(blank)
    s2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text_frame.text = "二枚目のスライド。"
    prs.save(os.path.join(OUT, "sample.pptx"))


def make_large_docx(n_para=50000):
    import docx
    d = docx.Document()
    for i in range(n_para):
        d.add_paragraph(f"Paragraph {i}: the quick brown fox jumps over the lazy dog.")
    d.save(os.path.join(OUT, "big.docx"))


def make_large_xlsx(n_row=100000, n_col=10):
    import openpyxl
    wb = openpyxl.Workbook(write_only=True)
    ws = wb.create_sheet("Big")
    for r in range(1, n_row + 1):
        ws.append([f"r{r}c{c}" if c % 2 == 0 else r * 1000 + c for c in range(n_col)])
    wb.save(os.path.join(OUT, "big.xlsx"))


def make_large_pptx(n_slide=1000):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(n_slide):
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = (
            f"Slide {i} body text content for benchmarking.")
    prs.save(os.path.join(OUT, "big.pptx"))


# ---------------------------------------------------------------------------
# Legacy binary samples (LibreOffice UNO)
# ---------------------------------------------------------------------------

def _connect():
    import uno
    ctx = uno.getComponentContext()
    resolver = ctx.ServiceManager.createInstanceWithContext(
        "com.sun.star.bridge.UnoUrlResolver", ctx)
    remote = resolver.resolve(
        "uno:socket,host=localhost,port=2003;urp;StarOffice.ComponentContext")
    return remote.ServiceManager.createInstanceWithContext(
        "com.sun.star.frame.Desktop", remote)


def _save(doc, name, filt):
    from com.sun.star.beans import PropertyValue
    pv = PropertyValue()
    pv.Name, pv.Value = "FilterName", filt
    doc.storeToURL("file:///" + OUT.replace("\\", "/") + "/" + name, (pv,))


def _add_textbox(doc, page, x, y, w, h, text):
    import uno
    shape = doc.createInstance("com.sun.star.drawing.TextShape")
    size = uno.createUnoStruct("com.sun.star.awt.Size")
    size.Width, size.Height = w, h
    pos = uno.createUnoStruct("com.sun.star.awt.Point")
    pos.X, pos.Y = x, y
    page.add(shape)
    shape.Size, shape.Position = size, pos
    shape.setString(text)


def shapes_calc(desktop):
    calc = desktop.loadComponentFromURL("private:factory/scalc", "_blank", 0, ())
    sheet = calc.Sheets.getByIndex(0)
    sheet.getCellByPosition(0, 0).setString("CellText")
    _add_textbox(calc, sheet.DrawPage, 1000, 1000, 6000, 2000, "Shape in Excel 図形テキスト")
    cell = sheet.getCellByPosition(2, 0)
    cell.setString("Commented")
    cell.Annotation.setString("This is a comment コメント")
    _save(calc, "shapes.xls", "MS Excel 97")


def shapes_writer(desktop):
    import uno
    writer = desktop.loadComponentFromURL("private:factory/swriter", "_blank", 0, ())
    writer.Text.setString("Body paragraph 本文です。")
    frame = writer.createInstance("com.sun.star.text.TextFrame")
    size = uno.createUnoStruct("com.sun.star.awt.Size")
    size.Width, size.Height = 6000, 2000
    frame.setPropertyValue("Width", 6000)
    frame.setPropertyValue("Height", 2000)
    writer.Text.insertTextContent(writer.Text.createTextCursor(), frame, False)
    frame.getText().setString("Shape in Word 図形の中のテキスト")
    _save(writer, "shapes.doc", "MS Word 97")


def shapes_impress(desktop):
    from com.sun.star.beans import PropertyValue
    hidden = PropertyValue()
    hidden.Name, hidden.Value = "Hidden", True
    impress = desktop.loadComponentFromURL("private:factory/simpress", "_blank", 0, (hidden,))
    slide = impress.DrawPages.getByIndex(0)
    if slide.Count > 0:
        slide.getByIndex(0).setString("Slide Title タイトル")
    _add_textbox(impress, slide, 2000, 5000, 15000, 3000, "Shape in PowerPoint 図形テキスト")
    _save(impress, "shapes.ppt", "MS PowerPoint 97")


def big_calc(desktop, n=30000):
    calc = desktop.loadComponentFromURL("private:factory/scalc", "_blank", 0, ())
    sheet = calc.Sheets.getByIndex(0)
    data = tuple(("MK%06d" % i, "和%06d" % i) for i in range(n))
    sheet.getCellRangeByPosition(0, 0, 1, n - 1).setDataArray(data)
    # A string longer than a BIFF8 record, to span Continue records.
    sheet.getCellByPosition(3, 0).setString("LONGSTART " + ("x" * 20000) + " LONGEND_MK")
    _save(calc, "big.xls", "MS Excel 97")


def big_writer(desktop, n=8000):
    from com.sun.star.text.ControlCharacter import PARAGRAPH_BREAK
    writer = desktop.loadComponentFromURL("private:factory/swriter", "_blank", 0, ())
    text = writer.Text
    cur = text.createTextCursor()
    for i in range(n):
        text.insertString(cur, "P%06d:MK%06d 日本語%06d" % (i, i, i), False)
        text.insertControlCharacter(cur, PARAGRAPH_BREAK, False)
    text.insertString(cur, "LONGSTART " + ("y" * 30000) + " LONGEND_MK", False)
    _save(writer, "big.doc", "MS Word 97")


def big_impress(desktop, n=250):
    from com.sun.star.beans import PropertyValue
    hidden = PropertyValue()
    hidden.Name, hidden.Value = "Hidden", True
    impress = desktop.loadComponentFromURL("private:factory/simpress", "_blank", 0, (hidden,))
    slides = impress.DrawPages
    while slides.Count < n:
        slides.insertNewByIndex(slides.Count)
    for i in range(n):
        _add_textbox(impress, slides.getByIndex(i), 1000, 1000, 15000, 3000,
                     "S%06d:MK%06d スライド%06d" % (i, i, i))
    _save(impress, "big.ppt", "MS PowerPoint 97")


def _run_uno(funcs):
    import traceback
    desktop = _connect()
    for fn in funcs:
        try:
            fn(desktop)
            print(fn.__name__, "OK")
        except Exception:
            print(fn.__name__, "FAILED")
            traceback.print_exc()


# ---------------------------------------------------------------------------

def main():
    cmd = sys.argv[1] if len(sys.argv) > 1 else "samples"
    if cmd == "samples":
        make_docx(); make_xlsx(); make_pptx()
        print("wrote sample.docx/.xlsx/.pptx to", OUT)
    elif cmd == "large":
        make_large_docx(); make_large_xlsx(); make_large_pptx()
        print("wrote big.docx/.xlsx/.pptx to", OUT)
    elif cmd == "shapes":
        _run_uno([shapes_calc, shapes_writer, shapes_impress])
    elif cmd == "big":
        _run_uno([big_calc, big_writer, big_impress])
    else:
        sys.exit("unknown command %r (use: samples | large | shapes | big)" % cmd)


if __name__ == "__main__":
    main()
